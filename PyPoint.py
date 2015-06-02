#!/usr/bin/python2.70
import os
import subprocess
import sys
import pygtk
pygtk.require('2.0')
from gi.repository import Gtk as gtk
from pptx import *
import string
import PIL
from PIL import Image
import xlrd


pagesList = ["welcomePage","titlePage","textPage","picPage","excelPage", "finishPage"]
pageInit = 0
prs = Presentation
class PageControl:

	# Variables created prior to any object creation
	window = None
	gladeFile = "MainGUI2.glade"
	builder = None
	fileChooser = None
	handlerObject = None
	fileLocation = None
	realFileLocation = None
	tableToggleButton = None
	chartToggleButton = None

	def __init__(self, page):

		# Variables created for each PageControl object

		self.page = page
		PageControl.builder = gtk.Builder()
		PageControl.builder.add_from_file(PageControl.gladeFile)
		PageControl.handlerObject = Handler(pagesList[page])
		PageControl.builder.connect_signals(PageControl.handlerObject)
		PageControl.window = PageControl.builder.get_object(pagesList[page])
		PageControl.window.show_all()


class Handler(PageControl):

	def __init__(self, currentPage):
		self.currentPage = currentPage

	def on_QuitButton_clicked(self, button):
		gtk.main_quit()

	def on_NextButton_clicked(self, button):
		for x in range(len(pagesList)):
			if pagesList[x] == self.currentPage:
				SlideCreator(x)
				PageControl.window.hide()
				if x != 5:
					PageControl(x+1)

	def on_FinishButton_clicked(self, button):
		PageControl.window.hide()
		subprocess.call(["xdg-open", PageControl.realFileLocation])
		gtk.main_quit()
	def on_currentFolderChanged(self, file_chooser):
		folderLocation = file_chooser.get_uri()
		realfolderLocation = str(folderLocation[7:]) + '/'
		return realfolderLocation

	def on_currentFileChange(self, file_chooser):
		fileLocation = file_chooser.get_filename()
		return fileLocation


	def on_CreateSlide_clicked(self, button):
		bulletBuffer = PageControl.builder.get_object("BulletTextView").get_buffer()
		bulletText = bulletBuffer.get_text(bulletBuffer.get_start_iter(), bulletBuffer.get_end_iter(), True)
		bulletButton = PageControl.builder.get_object("BulletCreateSlideButton")
		bulletButton.set_name("bulletButton")
		numberBuffer = PageControl.builder.get_object("NumberTextView").get_buffer()
		numberText = numberBuffer.get_text(numberBuffer.get_start_iter(), numberBuffer.get_end_iter(), True)
		numberButton = PageControl.builder.get_object("NumberCreateSlideButton")
		numberButton.set_name("numberButton")
		paragraphBuffer = PageControl.builder.get_object("ParagraphTextView").get_buffer()
		paragraphText = paragraphBuffer.get_text(paragraphBuffer.get_start_iter(), paragraphBuffer.get_end_iter(), True)
		paragraphButton = PageControl.builder.get_object("ParagraphCreateSlideButton")
		paragraphButton.set_name("paragraphButton")
		
		if button.get_name() == "bulletButton":
			bulletSlide = prs.slides.add_slide(prs.slide_layouts[1])
			listofLines = bulletText.split('\n')
			body_shape = bulletSlide.placeholders[1]
			tf = body_shape.text_frame
			tf.text = listofLines[0]
			for x in range(1, len(listofLines)):
				p = tf.add_paragraph()
				p.text = listofLines[x]
				p.level = 0
			bulletBuffer.set_text('')
		if button.get_name() == "numberButton":
			numberSlide = prs.slides.add_slide(prs.slide_layouts[1])
			body_shape = numberSlide.placeholders[1]
			tf = body_shape.text_frame
			tf.text = numberText
			numberBuffer.set_text('')
		if button.get_name() == "paragraphButton":
			paragraphSlide = prs.slides.add_slide(prs.slide_layouts[1])
			body_shape = paragraphSlide.placeholders[1]
			body_shape.text = paragraphText
			paragraphBuffer.set_text('')
		prs.save(PageControl.realFileLocation)

	def on_pictureFileActivated(self, file_chooser):
		picFileLocation = file_chooser.get_filename()
		picFormat = picFileLocation[len(picFileLocation) - 4:]
		picSlide = prs.slides.add_slide(prs.slide_layouts[6])
		img = Image.open(picFileLocation)
		if img.size[1] >= img.size[0] * 1.5:
			basewidth = 200
		else:
			basewidth = 400
		wpercent = (basewidth/float(img.size[0]))
		hsize = int((float(img.size[1])*float(wpercent)))
		img = img.resize((basewidth,hsize), Image.ANTIALIAS)
		tempPath = file_chooser.get_current_folder() + '/temp_resized' + picFormat
		img.save(tempPath)
		left = Inches(2)
		top = Inches(1)
		pic = picSlide.shapes.add_picture(tempPath, left, top)
		prs.save(PageControl.realFileLocation)
		os.remove(tempPath)

	def if_Toggled(self, checkButton):
		chartCheckButton = PageControl.builder.get_object("ChartCheckButton")
		chartCheckButton.set_name("ChartCheckButton")
		tableCheckButton = PageControl.builder.get_object("TableCheckButton")
		tableCheckButton.set_name("TableCheckButton")
		if checkButton.get_name() == "ChartCheckButton":
			if PageControl.chartToggleButton == None:
				print "Chart is On"
				PageControl.chartToggleButton = "Chart"
			elif PageControl.chartToggleButton == "Chart":
				print "Chart is Off"
				PageControl.chartToggleButton = None
		if checkButton.get_name() == "TableCheckButton":
			if PageControl.tableToggleButton == None:
				print "Table is On"
				PageControl.tableToggleButton = "Table"
			elif PageControl.tableToggleButton == "Table":
				print "Table is Off"
				PageControl.tableToggleButton = None


class SlideCreator(PageControl):
	folderLocation = None

	def welcomePage(self):

		fileChooser = PageControl.builder.get_object("fileChooser")
		SlideCreator.folderLocation = PageControl.handlerObject.on_currentFolderChanged(fileChooser)

	def titlePage(self):

		global prs
		TextBox1 = PageControl.builder.get_object("TitleTextBox").get_text()
		PageControl.realFileLocation = SlideCreator.folderLocation + TextBox1 + '.pptx'
		tempPrs = Presentation()
		tempPrs.save(PageControl.realFileLocation)
		prs = tempPrs
		for slide in prs.slides:
			print slide
		rawGroupMemberText = PageControl.builder.get_object("GroupMembersTextBox").get_text()
		newText = rawGroupMemberText.replace(', ', '\n')
		titleSlide = prs.slides.add_slide(prs.slide_layouts[0])
		title = titleSlide.shapes.title
		subtitle = titleSlide.placeholders[1]
		title.text = TextBox1
		subtitle.text = newText
		prs.save(PageControl.realFileLocation)

	def textPage(self):
		pass
	def picPage(self):
		pass
	def excelPage(self):

		sheetValueTextBox = PageControl.builder.get_object("ExcelSheetValue")	
		fileChooser = PageControl.builder.get_object("fileChooserExcel")

		def addChart(prs,filename,sheetNum):	#adds charts
			#create a blank chart slide
			slide = prs.slides.add_slide(prs.slide_layouts[5])
			#open the worksheet
			f=xlrd.open_workbook(filename)
			sheet=f.sheet_by_index(sheetNum-1)
			#get the data
			chart_data = ChartData()
			categories=[]
			label = ''
			for col in range(0,sheet.ncols):
				series=[]	#list to represent the data in that row
		#		print "Col",col
				for row in range(0,sheet.nrows):
		#			print "row",row
					if col == 0 and (sheet.cell(row,col) is not xlrd.empty_cell):	#for every row, if we are looking at the first column (index 0), add that to the list of categories
						cat = sheet.cell_value(row,col)
						cat = cat.encode('ascii','ignore')
		#				print "adding cat",cat
						categories.append(cat)
					elif row == 0:
						label = sheet.cell_value(row,col).encode('ascii','ignore')
					elif row >= 1:
		#				print "adding value",sheet.cell_value(col,row)
						series.append(sheet.cell_value(row,col))
				if len(series) > 0 and len(label) > 0:
					chart_data.add_series(label,series) #add the data for the row to the ChartData object
			categories.remove('')
			chart_data.categories=categories	#add the categories to the ChartData object
			#write it into a slide object
			x,y,cx,cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
			chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,x,y,cx,cy,chart_data).chart
			#add a legend
			chart.has_legend = True
			chart.legend.position=XL_LEGEND_POSITION.RIGHT
			chart.legend.include_in_layout = False

			prs.save(PageControl.realFileLocation)
			
			#Program parameters:
				#sheetNum: must be input as a number, that way sheet number corresponds to sheet index 0
				#see if loops for chartType parameters

		def addTable(prs,filename,sheetNum):	#adds tables
			#create blank table slide
			slide=prs.slides.add_slide(prs.slide_layouts[5])
			shapes = slide.shapes
			#open the worksheet
			f = xlrd.open_workbook(filename)
			sheet = f.sheet_by_index(sheetNum-1)
			#create the blank table
			rows = sheet.nrows
			cols = sheet.ncols
			left = top = Inches(2.0)
			width = Inches(2.0)
			height = Inches(0.8)
			table = shapes.add_table(rows,cols,left,top,width,height).table
			#get the data and write into the slide
			for row in range (0, sheet.nrows):
				for col in range (0, sheet.ncols):
					table.cell(row,col).text = str(sheet.cell_value(row,col))

			prs.save(PageControl.realFileLocation)

		if PageControl.tableToggleButton == "Table":
			print "Adding Table..."
			addTable(prs, PageControl.handlerObject.on_currentFileChange(fileChooser), int(sheetValueTextBox.get_text()))
			print "Table Added"
		if PageControl.chartToggleButton == "Chart":
			print "Adding Chart..."
			addChart(prs, PageControl.handlerObject.on_currentFileChange(fileChooser), int(sheetValueTextBox.get_text()))
			print "Chart Added"

	def finishPage(self):
		pass

	def __init__(self, page):
		if page == 0:
			self.welcomePage()
		elif page == 1:
			self.titlePage()
		elif page == 2:
			self.textPage()
		elif page == 3:
			self.picPage()
		elif page == 4:
			self.excelPage()
		elif page == 5:
			self.finishPage()


if __name__ == "__main__":
	prezi = PageControl(pageInit)
	gtk.main()
